import os
import sys
import math
import warnings
warnings.filterwarnings("ignore")

# --- Dependency Check ---
REQUIRED_PACKAGES = {
    'pandas': 'pandas',
    'geopandas': 'geopandas',
    'shapely': 'shapely',
    'matplotlib': 'matplotlib',
    'contextily': 'contextily',
    'pptx': 'python-pptx',
    'shapefile': 'pyshp',
    'pyproj': 'pyproj',
    'openpyxl': 'openpyxl',
}

MISSING_DEPS = []
for module_name, pip_name in REQUIRED_PACKAGES.items():
    try:
        __import__(module_name)
    except ImportError:
        MISSING_DEPS.append(pip_name)

if MISSING_DEPS:
    print("=" * 60)
    print("ERROR: Missing required Python packages!")
    print("=" * 60)
    print(f"\nThe following packages are not installed:\n")
    for dep in MISSING_DEPS:
        print(f"  - {dep}")
    print(f"\nPlease run this command to install them:\n")
    print(f"  pip install {' '.join(MISSING_DEPS)}")
    print(f"\nIf geopandas/fiona fails, try:")
    print(f"  pip install geopandas --only-binary :all:")
    print("=" * 60)
    sys.exit(1)

import pandas as pd
import geopandas as gpd
from shapely.geometry import Polygon, Point
from shapely.ops import unary_union
import matplotlib
matplotlib.use('Agg')  # Force non-interactive backend for headless environments
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import contextily as cx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import shapefile

def safe_add_basemap(ax, crs, source, attribution=False):
    """Wrapper around cx.add_basemap that gracefully handles network failures."""
    try:
        cx.add_basemap(ax, crs=crs, source=source, attribution=attribution)
    except Exception as e:
        print(f"  [WARNING] Could not download map tiles: {e}")
        print(f"  [WARNING] Using solid background fallback. Check your internet connection.")
        # Fallback: just set a dark background so the plot still renders
        ax.set_facecolor('#2d2d2d')

from config import MR_DIR, SITES_CSV, AIRPORT_DIR, PROPOSALS_XLSX, OUT_DIR, EVIDENCE_DIR

SHP_PATH = os.path.join(AIRPORT_DIR, "airport_border.shp")

def get_sector_polygon(lon, lat, y_lat_for_scale, azimuth, radius_m=200, angle_deg=65):
    import pyproj
    from shapely.geometry import shape
    transformer_to_3857 = pyproj.Transformer.from_crs("EPSG:4326", "EPSG:3857", always_xy=True)
    x_3857, y_3857 = transformer_to_3857.transform(lon, lat)
    
    start_angle = math.radians(90 - azimuth - angle_deg/2)
    end_angle = math.radians(90 - azimuth + angle_deg/2)
    
    points = [(x_3857, y_3857)]
    for i in range(11):
        angle = start_angle + (end_angle - start_angle) * (i / 10.0)
        x_pt = x_3857 + radius_m * math.cos(angle)
        y_pt = y_3857 + radius_m * math.sin(angle)
        points.append((x_pt, y_pt))
    points.append((x_3857, y_3857))
    return Polygon(points)

def get_rsrp_color(val):
    if val < -115: return '#FF0000'
    elif val < -110: return '#FFC000'
    elif val < -105: return '#FFFF00'
    elif val < -95: return '#92D050'
    else: return '#00B050'

def main():
    if len(sys.argv) < 2:
        print("Usage: python generate_single_pptx.py <Airport Name>")
        return
    
    target_airport = sys.argv[1].strip()
    print(f"Generating PPTX for {target_airport}")
    
    # 1. Load Airport
    import pyproj
    transformer_to_3857 = pyproj.Transformer.from_crs("EPSG:4326", "EPSG:3857", always_xy=True)
    
    sf = shapefile.Reader(SHP_PATH)
    fields = sf.fields[1:]
    field_names = [field[0] for field in fields]
    
    apt = None
    for shape_rec in sf.iterShapeRecords():
        rec = dict(zip(field_names, shape_rec.record))
        name = rec.get('Airport', '').strip().replace('\\\\', '')
        if name.lower() == target_airport.lower():
            poly = Polygon(shape_rec.shape.points)
            gdf_3857 = gpd.GeoDataFrame({'geometry': [poly]}, crs="EPSG:3857")
            gdf_4326 = gdf_3857.to_crs(epsg=4326)
            
            # total_bounds of 4326 is minlon, minlat, maxlon, maxlat
            l_min, la_min, l_max, la_max = gdf_4326.total_bounds
            
            # Extract 3857 bounds
            xmin, ymin, xmax, ymax = gdf_3857.total_bounds
            
            apt = {
                'name': name,
                'gdf_3857': gdf_3857,
                'bounds': (l_min, la_min, l_max, la_max),
                'bounds_3857': (xmin, ymin, xmax, ymax)
            }
            break
            
    if not apt:
        print(f"Airport {target_airport} not found!")
        return
        
    airport_poly_3857 = apt['gdf_3857'].geometry.iloc[0]
    xmin, ymin, xmax, ymax = apt['bounds_3857']
    l_min, la_min, l_max, la_max = apt['bounds']
    
    # 2. Load Sites
    df_cells = pd.read_csv(SITES_CSV)
    df_cells['Longitude'] = pd.to_numeric(df_cells['Longitude'], errors='coerce')
    df_cells['Latitude'] = pd.to_numeric(df_cells['Latitude'], errors='coerce')
    mask_ex = (
        (df_cells['Longitude'] >= l_min) & (df_cells['Longitude'] <= l_max) &
        (df_cells['Latitude'] >= la_min) & (df_cells['Latitude'] <= la_max)
    )
    df_cells_apt = df_cells[mask_ex].dropna(subset=['Longitude', 'Latitude'])
    
    gdf_sectors = []
    for _, row in df_cells_apt.iterrows():
        poly = get_sector_polygon(row['Longitude'], row['Latitude'], row['Latitude'], float(row.get('Azimuth', 0)), radius_m=200, angle_deg=65)
        gdf_sectors.append(poly)
    gdf_sectors = gpd.GeoDataFrame(geometry=gdf_sectors, crs="EPSG:3857") if gdf_sectors else gpd.GeoDataFrame()

    # Load Proposals
    df_prop = pd.read_excel(PROPOSALS_XLSX)
    df_prop_apt = df_prop[df_prop['Airport'].astype(str).str.lower() == target_airport.lower()]
    
    global_additional = []
    global_newsite = []
    global_calc_sectors = []
    
    for _, row in df_prop_apt.iterrows():
        is_new = ("_ARPT_" in str(row['Site ID'])) or ("New Site" in str(row.get('Remark', '')))
        remark = str(row.get('Remark', ''))
        az = float(row.get('Azimuth', 0))
        rad = float(row.get('Radius', 600))
        if pd.isna(rad): rad = 600
        bw = 33 if 'Change Antenna' in remark else 65
        
        poly_viz = get_sector_polygon(row['Longitude'], row['Latitude'], row['Latitude'], az, radius_m=200, angle_deg=bw)
        poly_calc = get_sector_polygon(row['Longitude'], row['Latitude'], row['Latitude'], az, radius_m=rad, angle_deg=bw)
        
        if is_new:
            global_newsite.append(poly_viz)
        else:
            global_additional.append(poly_viz)
        global_calc_sectors.append(poly_calc)

    combined_calc_poly = unary_union(global_calc_sectors) if global_calc_sectors else None

    # 3. Load MR Data
    val_cols = {'RSRP': 'RSRP(All MRs) (dBm)', 'RSRQ': 'RSRQ(All MRs) (dB)'}
    processed_gdfs = {}
    
    for env in ['Combine', 'Indoor']:
        for source in ['MR', 'MDT']:
            for metric in ['RSRP', 'RSRQ']:
                if source == 'MR':
                    fname = f"{metric}_Airport_{env}.csv"
                else:
                    fname = f"{metric}_Airport_MDT_{env}.csv"
                
                path = os.path.join(MR_DIR, fname)
                if not os.path.exists(path): continue
                
                df = pd.read_csv(path)
                mask = (df['Longitude'] >= l_min) & (df['Longitude'] <= l_max) & (df['Latitude'] >= la_min) & (df['Latitude'] <= la_max)
                df_apt = df[mask].copy()
                
                if df_apt.empty: continue
                if len(df_apt) > 5000: df_apt = df_apt.sample(5000, random_state=42)
                
                if metric == 'RSRP':
                    df_apt['color'] = df_apt[val_cols['RSRP']].apply(get_rsrp_color)
                
                gdf = gpd.GeoDataFrame(df_apt, geometry=gpd.points_from_xy(df_apt.Longitude, df_apt.Latitude), crs="EPSG:4326")
                processed_gdfs[f"{metric}_{env}_{source}"] = gdf

    coverage_stats = {}
    mr_dfs = {k.split('_')[1]: processed_gdfs.get(k) for k in processed_gdfs if 'MR' in k and 'RSRP' in k}
    for dsrc, dname in [('Combine', 'Combine'), ('Indoor', 'Indoor')]:
        gdf_src = mr_dfs.get(dsrc)
        if gdf_src is None or gdf_src.empty: continue
        gdf_src_3857 = gdf_src.to_crs(epsg=3857)
        total_spots = len(gdf_src_3857[gdf_src_3857.geometry.within(airport_poly_3857)])
        if total_spots == 0: continue
        
        bad_spots_gdf = gdf_src_3857[(gdf_src_3857.geometry.within(airport_poly_3857)) & (gdf_src_3857[val_cols['RSRP']] < -105)]
        bad_count = len(bad_spots_gdf)
        before_pct = (total_spots - bad_count) / total_spots * 100
        after_pct = before_pct
        if combined_calc_poly and bad_count > 0:
            covered_bad = len(bad_spots_gdf[bad_spots_gdf.geometry.within(combined_calc_poly)])
            after_pct = (total_spots - bad_count + covered_bad) / total_spots * 100
        coverage_stats[dname] = {'before': before_pct, 'after': after_pct}
        
    # 4. Generate Slide Images (Removed Runway/Airport Border from UI)
    rsrq_cmap = plt.get_cmap('RdYlGn')
    rsrq_norm = plt.Normalize(vmin=-20, vmax=-5)
    
    img_paths = {}
    for metric in ['RSRP', 'RSRQ']:
        for env in ['Combine', 'Indoor']:
            num_plots = 4 if metric == 'RSRP' else 3
            fig, axes = plt.subplots(1, num_plots, figsize=(24 if metric == 'RSRP' else 18, 6.5), dpi=150, facecolor='#1F1F1F')
            plt.subplots_adjust(wspace=0.1, bottom=0.20)
            
            mr_key = f"{metric}_{env}_MR"
            mr_gdf = processed_gdfs.get(mr_key)
            mdt_gdf = processed_gdfs.get(f"{metric}_{env}_MDT")
            
            titles = ['Map', f'MR {metric} (Before)', f'MR {metric} (After)', f'MDT {metric}'] if metric == 'RSRP' else ['Map', f'MR {metric}', f'MDT {metric}']
            
            for i, ax in enumerate(axes):
                ax.set_xlim(xmin, xmax)
                ax.set_ylim(ymin, ymax)
                ax.axis('off')
                ax.set_title(titles[i], loc='left', fontweight='bold', fontsize=16, color='white')
                
                if i == 0:
                    safe_add_basemap(ax, crs="EPSG:3857", source=cx.providers.Esri.WorldImagery)
                else:
                    safe_add_basemap(ax, crs="EPSG:3857", source=cx.providers.CartoDB.DarkMatter)
                    
                if metric == 'RSRP':
                    if i == 1 and mr_gdf is not None:
                        mr_gdf_3857 = mr_gdf.to_crs(epsg=3857)
                        mr_gdf_3857['geometry'] = mr_gdf_3857.geometry.buffer(25, cap_style=3)
                        mr_gdf_3857.plot(ax=ax, color=mr_gdf_3857['color'], alpha=1.0, edgecolor='none', zorder=1)
                    elif i == 2 and mr_gdf is not None:
                        mr_pred = mr_gdf.copy()
                        mr_pred_3857 = mr_pred.to_crs(epsg=3857)
                        if combined_calc_poly is not None:
                            mask_bad = mr_pred[val_cols['RSRP']] < -105
                            mask_covered = mr_pred_3857.geometry.within(combined_calc_poly)
                            mr_pred.loc[mask_bad & mask_covered, 'color'] = '#92D050'
                        mr_pred_3857['geometry'] = mr_pred_3857.geometry.buffer(25, cap_style=3)
                        mr_pred_3857.plot(ax=ax, color=mr_pred['color'], alpha=1.0, edgecolor='none', zorder=1)
                        
                        if len(global_calc_sectors) > 0:
                            pass # We don't plot the huge 600m dashed lines anymore
                            
                    elif i == 3 and mdt_gdf is not None:
                        mdt_gdf_3857 = mdt_gdf.to_crs(epsg=3857)
                        mdt_gdf_3857['geometry'] = mdt_gdf_3857.geometry.buffer(10, cap_style=3)
                        mdt_gdf_3857.plot(ax=ax, color=mdt_gdf_3857['color'], alpha=1.0, edgecolor='none', zorder=1)
                else:
                    if i == 1 and mr_gdf is not None:
                        mr_gdf_3857 = mr_gdf.to_crs(epsg=3857)
                        mr_gdf_3857['geometry'] = mr_gdf_3857.geometry.buffer(25, cap_style=3)
                        mr_gdf_3857.plot(ax=ax, column=val_cols[metric], cmap=rsrq_cmap, norm=rsrq_norm, alpha=1.0, edgecolor='none', zorder=1)
                    elif i == 2 and mdt_gdf is not None:
                        mdt_gdf_3857 = mdt_gdf.to_crs(epsg=3857)
                        mdt_gdf_3857['geometry'] = mdt_gdf_3857.geometry.buffer(10, cap_style=3)
                        mdt_gdf_3857.plot(ax=ax, column=val_cols[metric], cmap=rsrq_cmap, norm=rsrq_norm, alpha=1.0, edgecolor='none', zorder=1)
                
                # Plot Airport Polygon on top (zorder=2.5) with cyan outline
                if airport_poly_3857 is not None:
                    gpd.GeoDataFrame(geometry=[airport_poly_3857], crs="EPSG:3857").plot(ax=ax, facecolor='none', edgecolor='cyan', linewidth=1.0, zorder=2.5)
                
                if len(gdf_sectors) > 0:
                    gdf_sectors.plot(ax=ax, facecolor='cyan', edgecolor='black', alpha=1.0, linewidth=0.5, zorder=5)
                    
                if i == 0 or (metric == 'RSRP' and i == 2) or (metric == 'RSRQ' and i == 1):
                    if len(global_additional) > 0:
                        add_gdf = gpd.GeoDataFrame(geometry=global_additional, crs="EPSG:3857")
                        add_gdf.plot(ax=ax, facecolor='magenta', edgecolor='black', alpha=1.0, linewidth=0.8, zorder=5)
                    if len(global_newsite) > 0:
                        ns_gdf = gpd.GeoDataFrame(geometry=global_newsite, crs="EPSG:3857")
                        ns_gdf.plot(ax=ax, facecolor='black', edgecolor='white', alpha=1.0, linewidth=0.8, zorder=5)

            if metric == 'RSRP' and env in coverage_stats:
                stats = coverage_stats[env]
                table_ax = fig.add_axes([0.35, 0.10, 0.3, 0.08])
                table_ax.axis('off')
                table_data = [
                    ["Condition", "Before", "After"],
                    ["> -105 dBm", f"{stats['before']:.1f}%", f"{stats['after']:.1f}%"]
                ]
                tbl = table_ax.table(cellText=table_data, cellLoc='center', loc='center')
                tbl.auto_set_font_size(False)
                tbl.set_fontsize(12)
                for (r, c), cell in tbl.get_celld().items():
                    cell.set_text_props(color='white')
                    if r == 0:
                        cell.set_text_props(weight='bold', color='white')
                        cell.set_facecolor('#444444')
                    else:
                        cell.set_facecolor('#333333')
                    cell.set_edgecolor('#555555')

            if metric == 'RSRP':
                patches = [
                    mpatches.Patch(color='#FF0000', label='< -115'),
                    mpatches.Patch(color='#FFC000', label='-115 to -110'),
                    mpatches.Patch(color='#FFFF00', label='-110 to -105'),
                    mpatches.Patch(color='#92D050', label='-105 to -95'),
                    mpatches.Patch(color='#00B050', label='>= -95'),
                    mpatches.Patch(facecolor='cyan', edgecolor='black', label='Existing Site'),
                    mpatches.Patch(facecolor='magenta', edgecolor='black', label='Additional Sector / High Gain'),
                    mpatches.Patch(facecolor='black', edgecolor='white', label='New Site')
                ]
                legend = fig.legend(handles=patches, loc='lower center', ncol=len(patches), bbox_to_anchor=(0.5, 0.0), frameon=False, fontsize=10)
                for text in legend.get_texts():
                    text.set_color("white")
            else:
                cax = fig.add_axes([0.3, 0.05, 0.4, 0.03])
                sm = plt.cm.ScalarMappable(cmap=rsrq_cmap, norm=rsrq_norm)
                sm.set_array([])
                cb = fig.colorbar(sm, cax=cax, orientation='horizontal', label='RSRQ (dB)')
                cb.set_label('RSRQ (dB)', color='white')
                cb.ax.xaxis.set_tick_params(color='white')
                cb.outline.set_edgecolor('white')
                plt.setp(plt.getp(cb.ax.axes, 'xticklabels'), color='white')

            img_path = os.path.join(OUT_DIR, f"{target_airport}_{metric}_{env}_slide.png")
            plt.savefig(img_path, bbox_inches='tight', pad_inches=0.1)
            plt.close(fig)
            img_paths[f"{metric}_{env}"] = img_path

    # Evidence Plot
    dx = xmax - xmin
    dy = ymax - ymin
    pad_x = dx * 0.3
    pad_y = dy * 0.3
    
    fig_ev, ax_ev = plt.subplots(figsize=(14, 12), dpi=150)
    ax_ev.set_xlim(xmin - pad_x, xmax + pad_x)
    ax_ev.set_ylim(ymin - pad_y, ymax + pad_y)
    ax_ev.axis('off')
    ax_ev.set_title(f"Coverage Evidence - {target_airport}\nBlue = Proposed Sector Footprint | Red = Original Bad RSRP Spots", fontweight='bold', fontsize=14)
    safe_add_basemap(ax_ev, crs="EPSG:3857", source=cx.providers.CartoDB.DarkMatter)
    
    if 'RSRP_Combine_MR' in processed_gdfs:
        gdf_ev = processed_gdfs['RSRP_Combine_MR'].to_crs(epsg=3857)
        bad_spots_ev = gdf_ev[(gdf_ev.geometry.within(airport_poly_3857)) & (gdf_ev[val_cols['RSRP']] < -105)]
        bad_spots_ev.plot(ax=ax_ev, color='red', markersize=8, alpha=0.9, zorder=5, label=f'Bad RSRP Spots ({len(bad_spots_ev)})')
    
    if len(global_calc_sectors) > 0:
        calc_gdf = gpd.GeoDataFrame(geometry=global_calc_sectors, crs="EPSG:3857")
        calc_gdf.plot(ax=ax_ev, facecolor='dodgerblue', edgecolor='navy', alpha=0.25, linewidth=0.8, zorder=4, label=f'Sector Footprint ({len(calc_gdf)} sectors)')
    
    if len(global_newsite) > 0:
        gpd.GeoDataFrame(geometry=global_newsite, crs="EPSG:3857").plot(ax=ax_ev, facecolor='black', edgecolor='white', alpha=1.0, linewidth=1, zorder=10, label=f'New Site Sectors ({len(global_newsite)})')
    if len(global_additional) > 0:
        gpd.GeoDataFrame(geometry=global_additional, crs="EPSG:3857").plot(ax=ax_ev, facecolor='magenta', edgecolor='black', alpha=1.0, linewidth=1, zorder=10, label=f'Additional Sectors ({len(global_additional)})')
    if len(gdf_sectors) > 0:
        gdf_sectors.plot(ax=ax_ev, facecolor='cyan', edgecolor='black', alpha=1.0, linewidth=0.5, zorder=10, label=f'Existing Sectors ({len(gdf_sectors)})')
        
    ax_ev.legend(loc='lower left', fontsize=10, framealpha=0.9)
    evidence_path = os.path.join(EVIDENCE_DIR, f"{target_airport}_Coverage_Evidence.png")
    os.makedirs(EVIDENCE_DIR, exist_ok=True)
    plt.savefig(evidence_path, bbox_inches='tight')
    plt.close(fig_ev)

    # 5. Compile PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slides_order = [
        ('RSRP', 'Combine', "RSRP Outdoor + Indoor", img_paths.get('RSRP_Combine')),
        ('RSRQ', 'Combine', "RSRQ Outdoor + Indoor", img_paths.get('RSRQ_Combine')),
        ('RSRP', 'Indoor', "RSRP Indoor Only", img_paths.get('RSRP_Indoor')),
        ('RSRQ', 'Indoor', "RSRQ Indoor Only", img_paths.get('RSRQ_Indoor')),
    ]
    
    # Slide 0: Cover Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 31, 31)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"{target_airport} Airport Coverage Improvement"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.name = "Courier New"
    
    import datetime
    today = datetime.datetime.now()
    def get_ordinal(n):
        if 11 <= (n % 100) <= 13: return str(n) + 'th'
        return str(n) + {1: 'st', 2: 'nd', 3: 'rd'}.get(n % 10, 'th')
    date_str = f"{get_ordinal(today.day)} {today.strftime('%B %Y')}"
    subtitle.text = date_str
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle.text_frame.paragraphs[0].font.name = "Courier New"

    for metric, env, title_str, img_path in slides_order:
        if not img_path or not os.path.exists(img_path): continue
        slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background to dark
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 31, 31)
        
        shapes = slide.shapes
        # Add a sleek custom title box
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{target_airport} - {title_str}"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        if metric == 'RSRQ':
            pic = slide.shapes.add_picture(img_path, Inches(1.0), Inches(1), width=Inches(11.333))
        else:
            pic = slide.shapes.add_picture(img_path, Inches(0.2), Inches(1.0), width=Inches(12.933))
            
    out_pptx = os.path.join(OUT_DIR, f"{target_airport}_Airport_Improvement.pptx")
    prs.save(out_pptx)
    print(f"Successfully generated {out_pptx}")

    # Delete the extra PNG files
    for path in img_paths.values():
        if os.path.exists(path):
            try: os.remove(path)
            except: pass

if __name__ == '__main__':
    main()
