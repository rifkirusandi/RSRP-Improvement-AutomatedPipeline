import os
import math
import numpy as np
import pandas as pd
import geopandas as gpd
from openpyxl.styles import Alignment, Font
from shapely.geometry import Polygon, Point
from shapely.ops import nearest_points
from sklearn.cluster import DBSCAN
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import gc
import matplotlib.colors as mcolors
import contextily as cx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import shapefile
from datetime import datetime
import pyproj
import fiona
from shapely.geometry import shape
import warnings
warnings.filterwarnings("ignore")

# Constants
MR_DIR = r"C:\Request\Airport Improvement\MR AIRPORT"
SITES_CSV = r"C:\Request\Airport Improvement\sites covering airport in all huawei foot-print v1.csv"
SHP_PATH = r"C:\Request\Airport Improvement\Internasional Airport Border\Internasional Airport Border.shp"
TERRITORY_PATH = r"C:\Request\Airport Improvement\Territory\Territory IOH 202605 - May v3.shp"
CLUTTER_PATH = r"C:\Request\Airport Improvement\Clutter\Morphology Indonesia V4.TAB"
RUNWAY_SHP = r"C:\Request\Airport Improvement\Runway\airport runway line buffer 150 m.shp"
TLP_CSV = r"C:\Request\Airport Improvement\TLP\All TLP NationWide Feb 2024 Update.csv"
OUT_DIR = "Output"
os.makedirs(OUT_DIR, exist_ok=True)

CLUTTER_RADII = {
    'DENSE URBAN': 400,
    'SUB URBAN': 1000,
    'URBAN': 600,
    'RURAL': 1500
}
DEFAULT_RADIUS = 600

def get_ordinal_suffix(day):
    if 11 <= day <= 13: return 'th'
    last_digit = day % 10
    if last_digit == 1: return 'st'
    elif last_digit == 2: return 'nd'
    elif last_digit == 3: return 'rd'
    else: return 'th'

today = datetime.now()
day = today.day
suffix = get_ordinal_suffix(day)
month_year = today.strftime('%B %Y')

print("Loading Airport Shapefile...")
sf = shapefile.Reader(SHP_PATH)
fields = sf.fields[1:]
field_names = [field[0] for field in fields]

airports = []
target_bbox = None
for shape_rec in sf.iterShapeRecords():
    rec_dict = dict(zip(field_names, shape_rec.record))
    airport_name = rec_dict.get('Airport', 'Unknown').strip().replace('\\', '')
    
    poly = Polygon(shape_rec.shape.points)
    gdf = gpd.GeoDataFrame(geometry=[poly], crs="EPSG:3857")
    gdf_4326 = gdf.to_crs(epsg=4326)
    minx, miny, maxx, maxy = gdf_4326.total_bounds
    
    # Expand bbox slightly for territory intersection
    target_bbox = (minx - 0.5, miny - 0.5, maxx + 0.5, maxy + 0.5)
    
    airports.append({
        'name': airport_name,
        'gdf_3857': gdf,
        'gdf_4326': gdf_4326,
        'lon_min': minx - 0.005,
        'lon_max': maxx + 0.005,
        'lat_min': miny - 0.005,
        'lat_max': maxy + 0.005
    })

print(f"Found {len(airports)} airports matching criteria.")

global_all_proposals = []

print("Loading Territory Shapefile...")
sf_terr = shapefile.Reader(TERRITORY_PATH)
terr_polys = []
for shape_rec in sf_terr.iterShapeRecords():
    bbox = shape_rec.shape.bbox
    if target_bbox:
        if bbox[0] > target_bbox[2] or bbox[2] < target_bbox[0] or bbox[1] > target_bbox[3] or bbox[3] < target_bbox[1]:
            continue
    try:
        if shape_rec.shape.shapeType != shapefile.NULL:
            terr_polys.append(shape(shape_rec.shape.__geo_interface__))
    except:
        pass
gdf_territory = gpd.GeoDataFrame(geometry=terr_polys, crs="EPSG:4326")
gdf_territory_4326 = gdf_territory

print("Loading cell sites...")
df_cells_all = pd.read_csv(SITES_CSV)
gdf_cells_all_3857 = gpd.GeoDataFrame(
    df_cells_all, geometry=gpd.points_from_xy(df_cells_all.Longitude, df_cells_all.Latitude), crs="EPSG:4326"
).to_crs(epsg=3857)

print("Loading Runway shapefile...")
try:
    sf_r = shapefile.Reader(RUNWAY_SHP)
    r_fields = [field[0] for field in sf_r.fields[1:]]
    r_polys = []
    r_airports = []
    for shape_rec in sf_r.iterShapeRecords():
        r_rec = dict(zip(r_fields, shape_rec.record))
        r_name = r_rec.get('Airport', 'Unknown').strip().replace('\\', '')
        poly = Polygon(shape_rec.shape.points)
        r_polys.append(poly)
        r_airports.append(r_name)
    gdf_runways = gpd.GeoDataFrame({'Airport': r_airports, 'geometry': r_polys}, crs="EPSG:4326")
    gdf_runways_3857 = gdf_runways.to_crs(epsg=3857)
    # The shapefile is a 150m buffer on each side (300m total width).
    # Real runway width is ~45-60m. Shrink by applying -120m negative buffer.
    gdf_runways_3857['geometry'] = gdf_runways_3857.geometry.buffer(-120)
    gdf_runways_3857 = gdf_runways_3857[~gdf_runways_3857.is_empty]
    print(f"  Runway polygons shrunk to ~60m realistic width ({len(gdf_runways_3857)} remaining)")
except Exception as e:
    print(f"Error loading runway shapefile manually: {e}")
    gdf_runways_3857 = gpd.GeoDataFrame(columns=['Airport', 'geometry'])

print("Loading TLP dataset...")
try:
    df_tlp = pd.read_csv(TLP_CSV)
    df_tlp['Longitude'] = pd.to_numeric(df_tlp['Longitude'], errors='coerce')
    df_tlp['Latitude'] = pd.to_numeric(df_tlp['Latitude'], errors='coerce')
    df_tlp = df_tlp.dropna(subset=['Longitude', 'Latitude'])
    gdf_tlp = gpd.GeoDataFrame(
        df_tlp, geometry=gpd.points_from_xy(df_tlp.Longitude, df_tlp.Latitude), crs="EPSG:4326"
    ).to_crs(epsg=3857)
    
    pass
except Exception as e:
    print(f"Error loading TLP CSV: {e}")
    gdf_tlp = gpd.GeoDataFrame()

print("Loading all MR/MDT data into memory...")
val_cols = {
    'RSRP': 'RSRP(All MRs) (dBm)',
    'RSRQ': 'RSRQ(All MRs) (dB)'
}

raw_data = {}
for metric in ['RSRP', 'RSRQ']:
    for env in ['Combine', 'Indoor']:
        for src in ['MR', 'MDT']:
            fname = f"{metric}_Airport_{env}.csv" if src == 'MR' else f"{metric}_Airport_MDT_{env}.csv"
            fpath = os.path.join(MR_DIR, fname)
            if os.path.exists(fpath):
                raw_data[f"{metric}_{env}_{src}"] = pd.read_csv(fpath)

print("Loading Clutter globally into memory...")
try:
    with fiona.open(CLUTTER_PATH) as src:
        features = list(src)
        if features:
            geoms = [shape(f['geometry']) for f in features if f['geometry'] is not None]
            props = [f['properties'] for f in features if f['geometry'] is not None]
            global_clutter_gdf = gpd.GeoDataFrame(props, geometry=geoms, crs=src.crs)
            global_clutter_gdf = global_clutter_gdf.to_crs(epsg=4326)
        else:
            global_clutter_gdf = gpd.GeoDataFrame()
except Exception as e:
    print(f"Error loading global clutter: {e}")
    global_clutter_gdf = gpd.GeoDataFrame()

def snap_azimuth(azimuth):
    if pd.isna(azimuth): return 0
    az = int(round(azimuth / 5.0) * 5)
    return az % 360

def get_sector_polygon(cx_m, cy_m, lat, azimuth, radius_m=200, angle_deg=65):
    if pd.isna(azimuth): azimuth = 0
    r_units = radius_m / math.cos(math.radians(lat))
    start_angle = 90 - (azimuth + angle_deg/2)
    end_angle = 90 - (azimuth - angle_deg/2)
    points = [(cx_m, cy_m)]
    for i in range(11):
        angle = start_angle + (end_angle - start_angle) * i / 10.0
        rad = math.radians(angle)
        x = cx_m + r_units * math.cos(rad)
        y = cy_m + r_units * math.sin(rad)
        points.append((x, y))
    points.append((cx_m, cy_m))
    return Polygon(points)

def get_rsrp_color(val):
    if val < -115: return '#FF0000' # Red
    elif val < -110: return '#FFC000' # Orange
    elif val < -105: return '#FFFF00' # Yellow
    elif val < -95: return '#92D050' # Light Green
    else: return '#00B050' # Dark Green

rsrq_cmap = plt.cm.get_cmap('RdYlGn')
rsrq_min = float('inf')
rsrq_max = float('-inf')
for k, df in raw_data.items():
    if 'RSRQ' in k:
        if len(df) > 0:
            v_min = df[val_cols['RSRQ']].min()
            v_max = df[val_cols['RSRQ']].max()
            if not pd.isna(v_min): rsrq_min = min(rsrq_min, v_min)
            if not pd.isna(v_max): rsrq_max = max(rsrq_max, v_max)

if rsrq_min == float('inf'):
    rsrq_norm = mcolors.Normalize(vmin=-20, vmax=0)
else:
    rsrq_norm = mcolors.Normalize(vmin=rsrq_min, vmax=rsrq_max)

transformer_to_3857 = pyproj.Transformer.from_crs("EPSG:4326", "EPSG:3857", always_xy=True)
transformer_to_4326 = pyproj.Transformer.from_crs("EPSG:3857", "EPSG:4326", always_xy=True)

def calculate_bearing(lon1, lat1, lon2, lat2):
    lon1, lat1, lon2, lat2 = map(math.radians, [lon1, lat1, lon2, lat2])
    dlon = lon2 - lon1
    x = math.sin(dlon) * math.cos(lat2)
    y = math.cos(lat1) * math.sin(lat2) - (math.sin(lat1) * math.cos(lat2) * math.cos(dlon))
    initial_bearing = math.atan2(x, y)
    initial_bearing = math.degrees(initial_bearing)
    return snap_azimuth((initial_bearing + 360) % 360)

def is_valid_azimuth(new_az, existing_azs, min_gap=90):
    for az in existing_azs:
        diff = abs(new_az - az)
        diff = min(diff, 360 - diff)
        if diff < min_gap:
            return False
    return True

def get_best_azimuths(site_pt_3857, radius_m, existing_azimuths, bad_spots_3857, max_sectors=3):
    dists = bad_spots_3857.geometry.distance(site_pt_3857)
    nearby = bad_spots_3857[dists <= radius_m]
    azimuths = list(existing_azimuths)
    
    if len(nearby) > 0:
        dx = nearby.geometry.x - site_pt_3857.x
        dy = nearby.geometry.y - site_pt_3857.y
        bearings = (np.degrees(np.arctan2(dx, dy)) + 360) % 360
        
        while len(azimuths) < max_sectors:
            best_count = -1
            best_az = -1
            for a in range(0, 360, 5):
                if not is_valid_azimuth(a, azimuths, 90): continue
                diff = np.abs(bearings - a)
                diff = np.minimum(diff, 360 - diff)
                count = np.sum(diff <= 32.5)
                if count > best_count:
                    best_count = count
                    best_az = a
                    
            if best_count > 0:
                azimuths.append(best_az)
                diff = np.abs(bearings - best_az)
                diff = np.minimum(diff, 360 - diff)
                mask = diff > 32.5
                bearings = bearings[mask]
            else:
                break
                
    while len(azimuths) < max_sectors:
        for offset in [120, 240, 90, 180, 270]:
            candidate = snap_azimuth((azimuths[0] + offset) % 360)
            if is_valid_azimuth(candidate, azimuths, 90):
                azimuths.append(candidate)
                break
                
    return azimuths

for apt in airports:
    name = apt['name']
    if "Kulon Progo" not in name: continue
    
    
    print(f"\nProcessing {name}...")
    
    l_min, l_max = apt['lon_min'], apt['lon_max']
    la_min, la_max = apt['lat_min'], apt['lat_max']
    
    df_cells = df_cells_all[
        (df_cells_all['Longitude'] >= l_min) & (df_cells_all['Longitude'] <= l_max) &
        (df_cells_all['Latitude'] >= la_min) & (df_cells_all['Latitude'] <= la_max)
    ].copy()
    
    # Filter cell sites to current airport bounds with a small buffer (~10km)
    buffer_deg = 0.1
    df_cells_apt = df_cells[
        (df_cells['Longitude'] >= l_min - buffer_deg) & (df_cells['Longitude'] <= l_max + buffer_deg) &
        (df_cells['Latitude'] >= la_min - buffer_deg) & (df_cells['Latitude'] <= la_max + buffer_deg)
    ].copy()
    
    # Store site states
    # site_id -> {'x': x, 'y': y, 'lat': lat, 'azimuths': []}
    sites_info = {}
    airport_proposals = []
    new_site_count = 1
    
    gdf_cells_3857 = gpd.GeoDataFrame()
    if len(df_cells_apt) > 0:
        gdf_cells = gpd.GeoDataFrame(
            df_cells_apt, geometry=gpd.points_from_xy(df_cells_apt.Longitude, df_cells_apt.Latitude), crs="EPSG:4326"
        )
        gdf_cells_3857 = gdf_cells.to_crs(epsg=3857)
        sectors = []
        for idx, row in gdf_cells_3857.iterrows():
            site_key = f"{row.geometry.x}_{row.geometry.y}"
            
            cell_name = str(row.get('Cell Name', '')).upper()
            if '_IN_I2_' in cell_name:
                c_type = 'IBC2M'
            elif '_IN_' in cell_name:
                c_type = 'IBS'
            else:
                c_type = 'MACRO'
                
            if site_key not in sites_info:
                sites_info[site_key] = {
                    'x': row.geometry.x, 
                    'y': row.geometry.y, 
                    'lat': row['Latitude'], 
                    'site_id': row.get('Site ID', 'EXISTING_SITE'), 
                    'azimuths': [],
                    'type': c_type
                }
            else:
                # Upgrade type if needed (IBC2M > MACRO > IBS)
                current = sites_info[site_key]['type']
                if c_type == 'IBC2M':
                    sites_info[site_key]['type'] = 'IBC2M'
                elif c_type == 'MACRO' and current == 'IBS':
                    sites_info[site_key]['type'] = 'MACRO'
                    
            if not pd.isna(row['Azimuth']):
                sites_info[site_key]['azimuths'].append(row['Azimuth'])
            sectors.append(get_sector_polygon(row.geometry.x, row.geometry.y, row['Latitude'], row['Azimuth'], radius_m=200))
        gdf_sectors = gpd.GeoDataFrame(geometry=sectors, crs="EPSG:3857")
    else:
        gdf_sectors = gpd.GeoDataFrame(geometry=[], crs="EPSG:3857")

    print(" Fetching Clutter...")
    clutter_gdf = gpd.GeoDataFrame()
    try:
        if len(global_clutter_gdf) > 0:
            clutter_gdf = gpd.clip(global_clutter_gdf, (l_min, la_min, l_max, la_max)).copy()
            if not clutter_gdf.empty:
                clutter_gdf['geometry'] = clutter_gdf.geometry.buffer(0)
                clutter_gdf = clutter_gdf.to_crs(epsg=3857)
    except Exception as e:
        print(" Clutter loading error:", e)

    def get_clutter_info(pt_3857):
        if len(clutter_gdf) > 0:
            intersecting = clutter_gdf[clutter_gdf.contains(pt_3857)]
            if not intersecting.empty:
                morpho = str(intersecting.iloc[0].get('Morpho', '')).strip().upper()
                for key, radius in CLUTTER_RADII.items():
                    if key in morpho:
                        return radius, key
        return DEFAULT_RADIUS, 'UNKNOWN' 
        
    def is_inside_clutter(pt_3857):
        """Check if a point falls inside the clutter map (i.e. on land, not ocean)"""
        if len(clutter_gdf) > 0:
            # Use spatial index for fast lookup
            possible_matches_idx = list(clutter_gdf.sindex.query(pt_3857, predicate='intersects'))
            return len(possible_matches_idx) > 0
        return True  # If no clutter loaded, allow placement

    def check_tlp(pt_3857):
        return True
    
    def get_isd_min(pt_3857):
        if len(clutter_gdf) > 0:
            intersecting = clutter_gdf[clutter_gdf.contains(pt_3857)]
            if not intersecting.empty:
                morpho = str(intersecting.iloc[0].get('Morpho', '')).strip().upper()
                if 'DENSE URBAN' in morpho: return 500
                elif 'SUB URBAN' in morpho: return 800
                elif 'URBAN' in morpho: return 600
                elif 'RURAL' in morpho: return 1200
        return 800 # default fallback
        
    processed_gdfs = {}
    
    # === OPTIMIZATION (Indoor MR RSRP) ===
    global_additional = []
    global_change_antenna = []
    global_newsite = []
    global_calc_sectors = []

    try:
        all_rsrp_dfs = [raw_data[k] for k in raw_data if 'RSRP' in k]
        if all_rsrp_dfs:
            df_in = pd.concat(all_rsrp_dfs, ignore_index=True)
            df_in = df_in[(df_in['Longitude'] >= l_min) & (df_in['Longitude'] <= l_max) &
                          (df_in['Latitude'] >= la_min) & (df_in['Latitude'] <= la_max)].copy()
            if len(df_in) > 0:
                gdf_in = gpd.GeoDataFrame(df_in, geometry=gpd.points_from_xy(df_in.Longitude, df_in.Latitude), crs="EPSG:4326")
                airport_poly_4326 = apt['gdf_4326'].geometry.iloc[0]
                airport_poly_3857 = apt['gdf_3857'].geometry.iloc[0]
                airport_exterior = airport_poly_3857.exterior
                
                # Calculation happens inside buffered polygon
                airport_poly_3857_buffered = airport_poly_3857.buffer(300)
                airport_poly_4326_buffered = gpd.GeoSeries([airport_poly_3857_buffered], crs="EPSG:3857").to_crs(epsg=4326).iloc[0]
                pts_inside = gdf_in[gdf_in.geometry.within(airport_poly_4326_buffered)]
                bad_spots = pts_inside[pts_inside['RSRP(All MRs) (dBm)'] < -105].copy()
                total_pts_in_airport = len(pts_inside)
                
                if total_pts_in_airport > 0 and len(bad_spots) > 0:
                    bad_spots_3857 = bad_spots.to_crs(epsg=3857)
                    uncovered_bad_spots = bad_spots_3857.copy()
                    
                    current_sites_info = {k: dict(v) for k, v in sites_info.items()}
                    for k in current_sites_info:
                        current_sites_info[k]['azimuths'] = list(current_sites_info[k]['azimuths'])
                        
                    # Iterative Simulation (100% coverage)
                    iteration = 0
                    while len(uncovered_bad_spots) > 0 and iteration < 100:
                        previous_uncovered_len = len(uncovered_bad_spots)
                        print(f"  Iteration {iteration}, Remaining Bad Spots: {len(uncovered_bad_spots)}")
                        iteration += 1
                        coords = np.array([(geom.x, geom.y) for geom in uncovered_bad_spots.geometry])
                        db = DBSCAN(eps=150, min_samples=3).fit(coords)
                        uncovered_bad_spots['cluster'] = db.labels_
                        
                        cluster_counts = uncovered_bad_spots[uncovered_bad_spots['cluster'] != -1]['cluster'].value_counts()
                        if len(cluster_counts) == 0:
                            break
                            
                        largest_cluster = cluster_counts.idxmax()
                        cluster_pts = uncovered_bad_spots[uncovered_bad_spots['cluster'] == largest_cluster]
                        
                        cx_c, cy_c = cluster_pts.geometry.x.mean(), cluster_pts.geometry.y.mean()
                        centroid_3857 = Point(cx_c, cy_c)
                        lon_c, lat_c = transformer_to_4326.transform(cx_c, cy_c)
                        
                        radius_m, morpho = get_clutter_info(centroid_3857)
                        sector_poly_calc = None
                        
                        placed = False
                        
                        # Step 1: Check existing sites for upgrades
                        closest_site_key = None
                        closest_dist = float('inf')
                        
                        for s_key, s_info in current_sites_info.items():
                            site_pt = Point(s_info['x'], s_info['y'])
                            dist = site_pt.distance(centroid_3857)
                            # We can consider sites within 1.2 * radius_m
                            if dist < closest_dist and dist <= radius_m * 1.2:
                                closest_dist = dist
                                closest_site_key = s_key
                                
                        if closest_site_key is not None:
                            s_info = current_sites_info[closest_site_key]
                            site_type = s_info.get('type', 'MACRO')
                            lon_s, lat_s = transformer_to_4326.transform(s_info['x'], s_info['y'])
                            bearing_to_cluster = calculate_bearing(lon_s, lat_s, lon_c, lat_c)
                            
                            # 1A: Propose IBC2M (If IBS only)
                            if site_type == 'IBS' and closest_dist <= radius_m:
                                if is_valid_azimuth(bearing_to_cluster, s_info['azimuths'], 90):
                                    s_info['azimuths'].append(bearing_to_cluster)
                                    s_info['type'] = 'IBC2M' # upgrade it
                                    airport_proposals.append({
                                        'Site ID': s_info.get('site_id', 'EXISTING_SITE'),
                                        'Longitude': lon_s,
                                        'Latitude': lat_s,
                                        'Azimuth': bearing_to_cluster,
                                        'Clutter': morpho,
                                        'Radius': radius_m,
                                        'Remark': 'Additional IBC2M',
                                        'Tower Provider ID': 'N/A',
                                        'Tower Provider Name': 'N/A'
                                    })
                                    poly_calc = get_sector_polygon(s_info['x'], s_info['y'], s_info['lat'], bearing_to_cluster, radius_m=radius_m, angle_deg=65)
                                    poly_viz = get_sector_polygon(s_info['x'], s_info['y'], s_info['lat'], bearing_to_cluster, radius_m=200, angle_deg=65)
                                    global_additional.append(poly_viz)
                                    global_calc_sectors.append(poly_calc)
                                    uncovered_bad_spots = uncovered_bad_spots[~uncovered_bad_spots.geometry.within(poly_calc)]
                                    placed = True
                            
                            # 1B: Change Antenna (If MACRO/IBC2M)
                            if not placed and site_type in ['MACRO', 'IBC2M']:
                                # Check if any existing sector points to the cluster
                                for i_az, az in enumerate(list(s_info['azimuths'])):
                                    diff = abs(bearing_to_cluster - az)
                                    diff = min(diff, 360 - diff)
                                    if diff <= 16.5: # 33 degree beamwidth
                                        # It points to the cluster. Evaluate High Gain Antenna
                                        new_radius = radius_m * 1.2
                                        poly_calc = get_sector_polygon(s_info['x'], s_info['y'], s_info['lat'], az, radius_m=new_radius, angle_deg=33)
                                        # Does it actually cover any bad spots?
                                        covered = uncovered_bad_spots[uncovered_bad_spots.geometry.within(poly_calc)]
                                        if len(covered) > 0:
                                            airport_proposals.append({
                                                'Site ID': s_info.get('site_id', 'EXISTING_SITE'),
                                                'Longitude': lon_s,
                                                'Latitude': lat_s,
                                                'Azimuth': az,
                                                'Clutter': morpho,
                                                'Radius': new_radius,
                                                'Remark': 'Change Antenna',
                                                'Tower Provider ID': 'N/A',
                                                'Tower Provider Name': 'N/A'
                                            })
                                            poly_viz = get_sector_polygon(s_info['x'], s_info['y'], s_info['lat'], az, radius_m=200, angle_deg=33)
                                            global_change_antenna.append(poly_viz) 
                                            global_calc_sectors.append(poly_calc)
                                            uncovered_bad_spots = uncovered_bad_spots[~uncovered_bad_spots.geometry.within(poly_calc)]
                                            placed = True
                                            break
                                        
                            # 1C: Additional Sector (If MACRO/IBC2M and Change Antenna didn't happen)
                            if not placed and site_type in ['MACRO', 'IBC2M'] and closest_dist <= radius_m:
                                max_sec = 3 if 'new_' in closest_site_key else 4
                                if len(s_info['azimuths']) < max_sec:
                                    if is_valid_azimuth(bearing_to_cluster, s_info['azimuths'], 90):
                                        s_info['azimuths'].append(bearing_to_cluster)
                                        airport_proposals.append({
                                            'Site ID': s_info.get('site_id', 'EXISTING_SITE'),
                                            'Longitude': lon_s,
                                            'Latitude': lat_s,
                                            'Azimuth': bearing_to_cluster,
                                            'Clutter': morpho,
                                            'Radius': radius_m,
                                            'Remark': 'Additional Sector',
                                            'Tower Provider ID': 'N/A',
                                            'Tower Provider Name': 'N/A'
                                        })
                                        poly_calc = get_sector_polygon(s_info['x'], s_info['y'], s_info['lat'], bearing_to_cluster, radius_m=radius_m, angle_deg=65)
                                        poly_viz = get_sector_polygon(s_info['x'], s_info['y'], s_info['lat'], bearing_to_cluster, radius_m=200, angle_deg=65)
                                        global_additional.append(poly_viz)
                                        global_calc_sectors.append(poly_calc)
                                        uncovered_bad_spots = uncovered_bad_spots[~uncovered_bad_spots.geometry.within(poly_calc)]
                                        placed = True
                                        
                        # Step 2: New Site
                        if not placed:
                            ns_pt = None
                            tlp_pid = 'N/A'
                            tlp_pname = 'N/A'
                            airport_runway = gdf_runways_3857[gdf_runways_3857['Airport'].str.contains(name, case=False, na=False)]
                            
                            def is_in_runway(pt):
                                if airport_runway.empty: return False
                                return any(airport_runway.geometry.contains(pt))
                                
                            if not gdf_tlp.empty:
                                tlp_dists = gdf_tlp.geometry.distance(centroid_3857)
                                nearby_tlp = gdf_tlp[tlp_dists <= 500].copy()
                                if not nearby_tlp.empty:
                                    nearby_tlp['dist'] = tlp_dists[tlp_dists <= 500]
                                    nearby_tlp = nearby_tlp.sort_values('dist')
                                    
                                    for _, tlp_row in nearby_tlp.iterrows():
                                        tlp_pt = tlp_row.geometry
                                        if not is_in_runway(tlp_pt) and is_inside_clutter(tlp_pt):
                                            # Ensure it's not within 50m of any existing airport site
                                            too_close = False
                                            if not gdf_cells_3857.empty:
                                                dists_to_ex = gdf_cells_3857.distance(tlp_pt)
                                                if dists_to_ex.min() <= 50:
                                                    too_close = True
                                            if too_close: continue
                                            ns_pt = tlp_pt
                                            tlp_pid = tlp_row.get('Tower Provider ID', 'N/A')
                                            tlp_pname = tlp_row.get('Tower Provider Name', 'N/A')
                                            break
                                        
                            if ns_pt is None:
                                if is_in_runway(centroid_3857):
                                    if not airport_runway.empty:
                                        geom = airport_runway.geometry.iloc[0].exterior
                                        p1, p2 = nearest_points(geom, centroid_3857)
                                        ns_pt = p1
                                else:
                                    ns_pt = centroid_3857
                            
                            # Enforce ISD against ALL sites (new + existing)
                            # Adaptive: reduce ISD progressively down to 500m minimum
                            # If still blocked, try shifting the site away from blocking site
                            if ns_pt is not None:
                                def find_isd_valid_point(candidate_pt):
                                    """Try to find a valid point respecting ISD. First try in-place with reducing ISD, then try shifting."""
                                    min_isd = get_isd_min(candidate_pt)
                                    
                                    # Phase 1: Try reducing ISD at the original location
                                    test_isd = min_isd
                                    while test_isd >= 500:
                                        blocked_by = None
                                        for s_key, s_info in current_sites_info.items():
                                            ex_pt = Point(s_info['x'], s_info['y'])
                                            d = candidate_pt.distance(ex_pt)
                                            if d < test_isd:
                                                blocked_by = (s_key, ex_pt, d)
                                                break
                                        if blocked_by is None:
                                            return candidate_pt  # Valid!
                                        test_isd -= 100
                                    
                                    # Phase 2: Shift away from the closest blocking site
                                    # Try 8 compass directions at the minimum ISD distance (500m)
                                    if blocked_by is not None:
                                        _, blocker_pt, _ = blocked_by
                                        import math
                                        for angle_offset in [0, 45, 90, 135, 180, 225, 270, 315]:
                                            # Direction from blocker to candidate, then offset
                                            dx = candidate_pt.x - blocker_pt.x
                                            dy = candidate_pt.y - blocker_pt.y
                                            base_angle = math.atan2(dy, dx)
                                            shift_angle = base_angle + math.radians(angle_offset)
                                            shift_dist = 550  # Just beyond 500m ISD
                                            
                                            shifted_pt = Point(
                                                blocker_pt.x + shift_dist * math.cos(shift_angle),
                                                blocker_pt.y + shift_dist * math.sin(shift_angle)
                                            )
                                            
                                            # Check this shifted point against ALL sites
                                            all_ok = True
                                            for s_key2, s_info2 in current_sites_info.items():
                                                ex_pt2 = Point(s_info2['x'], s_info2['y'])
                                                if shifted_pt.distance(ex_pt2) < 500:
                                                    all_ok = False
                                                    break
                                            
                                            # Also check it's still inside clutter/airport area
                                            if all_ok and is_inside_clutter(shifted_pt) and not is_in_runway(shifted_pt):
                                                print(f"DEBUG ISD: Shifted new site by {angle_offset}° to satisfy 500m ISD")
                                                return shifted_pt
                                    
                                    return None  # Could not find valid location
                                
                                ns_pt = find_isd_valid_point(ns_pt)
                                    
                            if ns_pt is not None:
                                lon_ns, lat_ns = transformer_to_4326.transform(ns_pt.x, ns_pt.y)
                                
                                if ns_pt.distance(centroid_3857) > 1:
                                    base_az = calculate_bearing(lon_ns, lat_ns, lon_c, lat_c)
                                else:
                                    base_az = 0
                                    
                                azs = [snap_azimuth(base_az), snap_azimuth(base_az + 120), snap_azimuth(base_az + 240)]
                                
                                new_site_key = f"new_{iteration}"
                                dummy_site_id = f"{name.upper().replace(' ', '_')}_ARPT_{new_site_count:03d}"
                                new_site_count += 1
                                current_sites_info[new_site_key] = {'x': ns_pt.x, 'y': ns_pt.y, 'lat': lat_ns, 'site_id': dummy_site_id, 'azimuths': azs, 'type': 'MACRO'}
                                
                                for az in azs:
                                    airport_proposals.append({
                                        'Site ID': dummy_site_id,
                                        'Longitude': lon_ns,
                                        'Latitude': lat_ns,
                                        'Azimuth': az,
                                        'Clutter': morpho,
                                        'Radius': radius_m,
                                        'Remark': 'New Site',
                                        'Tower Provider ID': tlp_pid,
                                        'Tower Provider Name': tlp_pname
                                    })
                                
                                for az in azs:
                                    poly_viz = get_sector_polygon(ns_pt.x, ns_pt.y, lat_ns, az, radius_m=200, angle_deg=65)
                                    poly_calc = get_sector_polygon(ns_pt.x, ns_pt.y, lat_ns, az, radius_m=radius_m, angle_deg=65)
                                    
                                    global_newsite.append(poly_viz)
                                    global_calc_sectors.append(poly_calc)
                                    
                                    uncovered_bad_spots = uncovered_bad_spots[~uncovered_bad_spots.geometry.within(poly_calc)]
                                placed = True
                        if not placed:
                            print(f"DEBUG: Dropping cluster {largest_cluster} because no point respects ISD, and no existing site can take a new sector.")
                            uncovered_bad_spots = uncovered_bad_spots[uncovered_bad_spots['cluster'] != largest_cluster]
                                    
                        if len(uncovered_bad_spots) == previous_uncovered_len:
                            print(f"DEBUG: No spots covered! Dropping cluster {largest_cluster} to avoid infinite loop.")
                            uncovered_bad_spots = uncovered_bad_spots[uncovered_bad_spots['cluster'] != largest_cluster]
    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"DEBUG EXCEPTION: {e}")

    opt_data = {
        'additional': gpd.GeoDataFrame(geometry=global_additional, crs="EPSG:3857"),
        'newsite': gpd.GeoDataFrame(geometry=global_newsite, crs="EPSG:3857")
    }

    for metric in ['RSRP', 'RSRQ']:
        for env in ['Combine', 'Indoor']:
            for src in ['MR', 'MDT']:
                key = f"{metric}_{env}_{src}"
                if key not in raw_data:
                    processed_gdfs[key] = None
                    continue
                df = raw_data[key]
                df_filtered = df[
                    (df['Longitude'] >= l_min) & (df['Longitude'] <= l_max) &
                    (df['Latitude'] >= la_min) & (df['Latitude'] <= la_max)
                ].copy()
                
                if len(df_filtered) == 0:
                    processed_gdfs[key] = None
                    continue
                
                gdf = gpd.GeoDataFrame(
                    df_filtered, geometry=gpd.points_from_xy(df_filtered.Longitude, df_filtered.Latitude), crs="EPSG:4326"
                )
                gdf_3857 = gdf.to_crs(epsg=3857)
                
                size_m = 20 if src == 'MDT' else 50
                avg_lat = gdf['Latitude'].mean()
                size_units = size_m / math.cos(math.radians(avg_lat))
                
                gdf_3857['geometry'] = gdf_3857.geometry.buffer(size_units/2, cap_style=3)
                
                if metric == 'RSRP':
                    gdf_3857['color'] = gdf_3857[val_cols['RSRP']].apply(get_rsrp_color)
                
                processed_gdfs[key] = gdf_3857

    xmin, ymin = transformer_to_3857.transform(l_min, la_min)
    xmax, ymax = transformer_to_3857.transform(l_max, la_max)
    
    # Calculate before/after coverage for Combine and Indoor
    coverage_stats = {}
    combined_calc_poly = None
    from shapely.ops import unary_union
    if global_calc_sectors:
        print(f"DEBUG: global_calc_sectors size: {len(global_calc_sectors)}")
        combined_calc_poly = unary_union(global_calc_sectors)
    
    airport_poly_3857 = apt['gdf_3857'].geometry.iloc[0]
    print(f"DEBUG: processed_gdfs keys: {list(processed_gdfs.keys())}")
    for k, v in processed_gdfs.items():
        if v is not None:
            print(f"DEBUG: {k} length: {len(v)}")
        else:
            print(f"DEBUG: {k} is None")
    
    mr_dfs = {k.split('_')[1]: processed_gdfs.get(k) for k in processed_gdfs if 'MR' in k and 'RSRP' in k}
    print(f"DEBUG: mr_dfs keys: {list(mr_dfs.keys())}")
    for dsrc, dname in [('Combine', 'Combine'), ('Indoor', 'Indoor')]:
        gdf_src = mr_dfs.get(dsrc)
        if gdf_src is None or gdf_src.empty: 
            print(f"DEBUG: {dsrc} is None or empty in mr_dfs")
            continue
        gdf_src_3857 = gdf_src.to_crs(epsg=3857)
        total_spots = len(gdf_src_3857[gdf_src_3857.geometry.within(airport_poly_3857)])
        print(f"DEBUG: {dsrc} total_spots within polygon: {total_spots}")
        if total_spots == 0: continue
        bad_spots_gdf = gdf_src_3857[(gdf_src_3857.geometry.within(airport_poly_3857)) & (gdf_src_3857[val_cols['RSRP']] < -105)]
        bad_count = len(bad_spots_gdf)
        before_pct = (total_spots - bad_count) / total_spots * 100
        after_pct = before_pct
        if combined_calc_poly and bad_count > 0:
            covered_bad = len(bad_spots_gdf[bad_spots_gdf.geometry.within(combined_calc_poly)])
            print(f"DEBUG: {dname} covered_bad: {covered_bad} out of {bad_count} bad spots")
            after_pct = (total_spots - bad_count + covered_bad) / total_spots * 100
        coverage_stats[dname] = {'before': before_pct, 'after': after_pct}
    print("Coverage Stats:", coverage_stats)

    # --- EVIDENCE PLOT ---
    if 'Combine' in mr_dfs and mr_dfs['Combine'] is not None and len(global_calc_sectors) > 0:
        # Add 30% padding around the airport bounds so nothing is cut off
        dx = xmax - xmin
        dy = ymax - ymin
        pad_x = dx * 0.3
        pad_y = dy * 0.3
        
        fig_ev, ax_ev = plt.subplots(figsize=(14, 12), dpi=150)
        ax_ev.set_xlim(xmin - pad_x, xmax + pad_x)
        ax_ev.set_ylim(ymin - pad_y, ymax + pad_y)
        ax_ev.axis('off')
        ax_ev.set_title(f"Coverage Evidence - {name}\nBlue = Proposed Sector Footprint | Red = Original Bad RSRP Spots", fontweight='bold', fontsize=14)
        cx.add_basemap(ax_ev, crs="EPSG:3857", source=cx.providers.CartoDB.Positron, attribution=False)
        
        # Plot airport boundary
        apt['gdf_3857'].plot(ax=ax_ev, facecolor='none', edgecolor='black', linewidth=2.5, linestyle='--', label='Airport Polygon')
        
        # Plot raw bad spots FIRST (underneath)
        gdf_ev = mr_dfs['Combine'].to_crs(epsg=3857)
        bad_spots_ev = gdf_ev[(gdf_ev.geometry.within(airport_poly_3857)) & (gdf_ev[val_cols['RSRP']] < -105)]
        bad_spots_ev.plot(ax=ax_ev, color='red', markersize=8, alpha=0.9, zorder=5, label=f'Bad RSRP Spots ({len(bad_spots_ev)})')
        
        # Plot the theoretical calculated footprints ON TOP
        calc_gdf = gpd.GeoDataFrame(geometry=global_calc_sectors, crs="EPSG:3857")
        calc_gdf.plot(ax=ax_ev, facecolor='dodgerblue', edgecolor='navy', alpha=0.25, linewidth=0.8, zorder=4, label=f'Sector Footprint ({len(calc_gdf)} sectors)')
        
        # Plot the proposed new site sectors
        if len(opt_data['newsite']) > 0:
            opt_data['newsite'].plot(ax=ax_ev, facecolor='purple', edgecolor='black', alpha=0.9, linewidth=1, zorder=6, label=f'New Site Sectors ({len(opt_data["newsite"])})')
        if len(opt_data['additional']) > 0:
            opt_data['additional'].plot(ax=ax_ev, facecolor='yellow', edgecolor='black', alpha=0.9, linewidth=1, zorder=6, label=f'Additional Sectors ({len(opt_data["additional"])})')
        
        # Plot existing site sectors
        if len(gdf_sectors) > 0:
            gdf_sectors.plot(ax=ax_ev, facecolor='orange', edgecolor='black', alpha=0.5, linewidth=0.5, zorder=3, label=f'Existing Sectors ({len(gdf_sectors)})')
            
        ax_ev.legend(loc='lower left', fontsize=10, framealpha=0.9)
        evidence_path = os.path.join("Evidence", f"{name}_Coverage_Evidence.png")
        plt.savefig(evidence_path, bbox_inches='tight')
        plt.close(fig_ev)
        gc.collect()
        print(f"Saved Evidence plot to {evidence_path}")
    # ---------------------


    img_paths = {}
    for metric in ['RSRP', 'RSRQ']:
        for env in ['Combine', 'Indoor']:
            num_plots = 4 if metric == 'RSRP' else 3
            fig, axes = plt.subplots(1, num_plots, figsize=(24 if metric == 'RSRP' else 18, 6.5), dpi=150, facecolor='#111111')
            plt.subplots_adjust(wspace=0.1, bottom=0.20)
            
            mr_key = f"{metric}_{env}_MR"
            mr_gdf = processed_gdfs.get(mr_key)
            mdt_gdf = processed_gdfs.get(f"{metric}_{env}_MDT")
            
            no_mr = True
            no_mdt = True
            if mr_gdf is not None and not mr_gdf.empty: no_mr = False
            if mdt_gdf is not None and not mdt_gdf.empty: no_mdt = False
            
            if no_mr and no_mdt:
                fig.text(0.5, 0.5, "MR & MDT not detected", fontsize=60, color='red', alpha=0.3, ha='center', va='center', rotation=30, fontweight='bold', zorder=100)
            elif no_mr:
                fig.text(0.5, 0.5, "MR not detected", fontsize=60, color='red', alpha=0.3, ha='center', va='center', rotation=30, fontweight='bold', zorder=100)
            elif no_mdt:
                fig.text(0.5, 0.5, "MDT not detected", fontsize=60, color='red', alpha=0.3, ha='center', va='center', rotation=30, fontweight='bold', zorder=100)
            
            titles = ['Map', f'MR {metric}', f'MR {metric} (After)', f'MDT {metric}'] if metric == 'RSRP' else ['Map', f'MR {metric}', f'MDT {metric}']
            
            for i, ax in enumerate(axes):
                ax.set_xlim(xmin, xmax)
                ax.set_ylim(ymin, ymax)
                ax.axis('off')
                ax.set_title(titles[i], loc='left', fontweight='bold', fontsize=16, color='white')
                
                if i == 0:
                    cx.add_basemap(ax, crs="EPSG:3857", source=cx.providers.Esri.WorldImagery, attribution=False)
                else:
                    cx.add_basemap(ax, crs="EPSG:3857", source=cx.providers.CartoDB.DarkMatter, attribution=False)
                    
                if metric == 'RSRP':
                    if i == 1 and mr_gdf is not None:
                        mr_gdf.plot(ax=ax, color=mr_gdf['color'], alpha=1.0, edgecolor='none')
                    elif i == 2 and mr_gdf is not None:
                        # Predictive MR (After)
                        mr_pred = mr_gdf.copy()
                        if combined_calc_poly is not None:
                            # Update colors for points falling within the new coverage
                            mask_bad = mr_pred[val_cols['RSRP']] < -105
                            mask_covered = mr_pred.geometry.within(combined_calc_poly)
                            mr_pred.loc[mask_bad & mask_covered, 'color'] = '#92D050' # Force them to Light Green (-105 to -95 range)
                        mr_pred.plot(ax=ax, color=mr_pred['color'], alpha=1.0, edgecolor='none')
                        
                        # (Removed sector radii overlay as per user request)
                        
                    elif i == 3 and mdt_gdf is not None:
                        mdt_gdf.plot(ax=ax, color=mdt_gdf['color'], alpha=1.0, edgecolor='none')
                else:
                    # RSRQ
                    if i == 1 and mr_gdf is not None:
                        mr_gdf.plot(ax=ax, column=val_cols[metric], cmap=rsrq_cmap, norm=rsrq_norm, alpha=1.0, edgecolor='none')
                    elif i == 2 and mdt_gdf is not None:
                        mdt_gdf.plot(ax=ax, column=val_cols[metric], cmap=rsrq_cmap, norm=rsrq_norm, alpha=1.0, edgecolor='none')
                
                if len(gdf_sectors) > 0:
                    gdf_sectors.plot(ax=ax, facecolor='orange', edgecolor='black', alpha=0.6, linewidth=0.5)
                    
                if (metric == 'RSRP' and (i == 1 or i == 2)) or (metric == 'RSRQ' and i == 1):
                    # In earlier version, opt_data was passed but we rewrote optimization logic so we must rebuild from airport_proposals
                    pass # Wait, we need to plot the new site locations and additional sectors
                    
                    if len(global_change_antenna) > 0:
                        ca_gdf = gpd.GeoDataFrame(geometry=global_change_antenna, crs="EPSG:3857")
                        ca_gdf.plot(ax=ax, facecolor='cyan', edgecolor='black', alpha=0.8, linewidth=0.8)
                    if len(global_additional) > 0:
                        add_gdf = gpd.GeoDataFrame(geometry=global_additional, crs="EPSG:3857")
                        add_gdf.plot(ax=ax, facecolor='yellow', edgecolor='black', alpha=0.8, linewidth=0.8)
                    if len(global_newsite) > 0:
                        ns_gdf = gpd.GeoDataFrame(geometry=global_newsite, crs="EPSG:3857")
                        ns_gdf.plot(ax=ax, facecolor='purple', edgecolor='black', alpha=0.8, linewidth=0.8)
                    
                if len(apt['gdf_3857']) > 0:
                    apt['gdf_3857'].plot(ax=ax, facecolor='none', edgecolor='cyan', linewidth=2)

            if metric == 'RSRP' and env in coverage_stats:
                stats = coverage_stats[env]
                table_ax = fig.add_axes([0.35, 0.10, 0.3, 0.08])
                table_ax.axis('off')
                table_data = [
                    ["Condition", "Before", "After"],
                    ["> -105 dBm", f"{stats['before']:.1f}%", f"{stats['after']:.1f}%"]
                ]
                tbl = table_ax.table(cellText=table_data, cellLoc='center', loc='center')
                tbl.auto_set_font_size(False)
                tbl.set_fontsize(12)
                for (r, c), cell in tbl.get_celld().items():
                    cell.set_text_props(color='white')
                    if r == 0:
                        cell.set_text_props(weight='bold', color='white')
                        cell.set_facecolor('#333333')
                    else:
                        cell.set_facecolor('#202020')
                    cell.set_edgecolor('#555555')

            if metric == 'RSRP':
                patches = [
                    mpatches.Patch(color='#FF0000', label='< -115'),
                    mpatches.Patch(color='#FFC000', label='-115 to -110'),
                    mpatches.Patch(color='#FFFF00', label='-110 to -105'),
                    mpatches.Patch(color='#92D050', label='-105 to -95'),
                    mpatches.Patch(color='#00B050', label='>= -95'),
                    mpatches.Patch(facecolor='#FFCC99', edgecolor='black', label='Existing Site'),
                    mpatches.Patch(facecolor='cyan', edgecolor='black', label='Change Antenna'),
                    mpatches.Patch(facecolor='yellow', edgecolor='black', label='Add Sector'),
                    mpatches.Patch(facecolor='purple', edgecolor='black', label='New Site')
                ]
                legend = fig.legend(handles=patches, loc='lower center', ncol=len(patches), bbox_to_anchor=(0.5, 0.0), frameon=False, fontsize=9)
                for text in legend.get_texts():
                    text.set_color("white")
            else:
                cax = fig.add_axes([0.3, 0.05, 0.4, 0.03])
                sm = plt.cm.ScalarMappable(cmap=rsrq_cmap, norm=rsrq_norm)
                sm.set_array([])
                cb = fig.colorbar(sm, cax=cax, orientation='horizontal', label='RSRQ (dB)')
                cb.set_label('RSRQ (dB)', color='white')
                cb.ax.xaxis.set_tick_params(color='white')
                cb.outline.set_edgecolor('white')
                plt.setp(plt.getp(cb.ax.axes, 'xticklabels'), color='white')

            img_path = os.path.join(OUT_DIR, f"{name}_{metric}_{env}_slide.png")
            plt.savefig(img_path, bbox_inches='tight', pad_inches=0.1)
            plt.close(fig)
            img_paths[f"{metric}_{env}"] = img_path

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank_layout)
    
    # Set background to dark gray
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(17, 17, 17)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = f"{name} Airport Coverage\nImprovement"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.name = 'Courier New'
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p2 = tf.add_paragraph()
    run_date = p2.add_run()
    run_date.text = f"\n{day}"
    run_date.font.size = Pt(24)
    run_date.font.name = 'Courier New'
    run_date.font.color.rgb = RGBColor(255, 255, 255)
    
    run_sup = p2.add_run()
    run_sup.text = suffix
    run_sup.font.size = Pt(24)
    run_sup.font.name = 'Courier New'
    run_sup.font.superscript = True
    run_sup.font.color.rgb = RGBColor(255, 255, 255)
    
    run_month = p2.add_run()
    run_month.text = f" {month_year}"
    run_month.font.size = Pt(24)
    run_month.font.name = 'Courier New'
    run_month.font.color.rgb = RGBColor(255, 255, 255)

    slides_order = [
        ('RSRP', 'Combine', f'{name} RSRP (Combine)'),
        ('RSRP', 'Indoor', f'{name} RSRP (Indoor)'),
        ('RSRQ', 'Combine', f'{name} RSRQ (Combine)'),
        ('RSRQ', 'Indoor', f'{name} RSRQ (Indoor)')
    ]

    for metric, env, title in slides_order:
        slide = prs.slides.add_slide(blank_layout)
        
        # Set background to dark gray
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(17, 17, 17)
        
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        img_path = img_paths.get(f"{metric}_{env}")
        if img_path and os.path.exists(img_path):
            # Center image vertically: image width=12.333, aspect ~3.7:1 -> height ~3.3"
            # Slide height=7.5, so center at ~(7.5 - 3.3)/2 = 2.1, but account for title
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(2.0), width=Inches(12.333))
            os.remove(img_path)

    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(4), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Thank You"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(0xd9, 0x00, 0x00)

    out_pptx = os.path.join(OUT_DIR, f"{name}_Airport_Improvement.pptx")
    prs.save(out_pptx)
    print(f"Saved {out_pptx}")
    if len(global_all_proposals) > 0:
        pd.DataFrame(global_all_proposals).to_excel(os.path.join(OUT_DIR, "All_Airports_Proposals.xlsx"), index=False)
    
    # --- ACCUMULATE EXCEL PROPOSALS ---
    if airport_proposals:
        global_all_proposals.extend(airport_proposals)

# --- EXPORT COMBINED EXCEL REPORT ---
if global_all_proposals:
    df_proposals = pd.DataFrame(global_all_proposals)
    excel_path = os.path.join(OUT_DIR, "All_Airports_Proposals.xlsx")
    with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
        df_proposals.to_excel(writer, index=False, sheet_name='Proposals')
        worksheet = writer.sheets['Proposals']
        
        # Center contents and bold headers
        for row in worksheet.iter_rows():
            for cell in row:
                cell.alignment = Alignment(horizontal='center')
        
        for cell in worksheet[1]:
            cell.font = Font(bold=True)
    print(f"Saved Combined Excel report to {excel_path}")

print("PPTX generated successfully!")
